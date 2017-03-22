from lxml import _elementpath as _dummy
from pptx import Presentation
import openpyxl
import pickle
import docx
from functools import*
import os
from tkinter import*
import collections.abc
import threading
import operator
from pdfparser import PDFParser, PDFDocument
from pdfinterp import PDFResourceManager, PDFPageInterpreter
from converter import PDFPageAggregator
from layout import LAParams, LTTextBox, LTTextLine
lkill=False
frame1=None
label2=None
myButton=None
includepdf=False
x={}
g=None
b=None
root1=None
root=None
box2=None
box4=None
box6=None
grint=False
gtxt=None
frameexists=False
upbox=False
labelenter=None
textbox=None
fram=None
frme=None
llll=None
llll2=None




def capspermutation(string):
   c=[]
   l=[]
   st3=string.upper()
   st2=string.lower()
   c.append(st2)
   c.append(st3)
   for alphabet in st2:
       l.append(alphabet)
   wrd=l[0].upper()
   e=1
   while e<len(l):
      wrd=wrd+l[e]
      e+=1
   c.append(wrd)
   return(c)

extlist=('.mp3','.ace','.doc','.java','xlsm','.mp4','.avi','.flv','.mkv','.jpg','.jpeg','.htm','.html','.pdf','.png','.exe','.txt','.docx','.xml','.gif','.odp','.pptx','.zip','.py','.cpp','.c','.rar','.xlsx')


def database(dir):
    global textbox
    try:
       textbox.insert(END,"{}\n".format(dir))
       textbox.see(END)
    except:
       print('no')
    b=os.listdir(dir)
    for entity in b:
        sleep(0.005)
        try:
          
           textbox.insert(END,"{}\n".format(entity))
        except:
           print('No')
        if entity=="$Recycle.Bin" or entity=="$RECYCLE.BIN" or entity=="System Volume Information" or "appdata" in entity.lower():
           continue
        if os.path.isdir(r'{}\{}'.format(dir,entity)):
           x[entity]=r'{}\{}'.format(dir,entity)
           try:
             database(r'{}\{}'.format(dir,entity))
           except:
             continue
        else:
           if '.mp3' in entity:
              if '.lnk' in entity or '.LNK' in entity:
                 continue                
              meta=getmeta(r'{}\{}'.format(dir,entity))
              tuple2=(entity,meta)
              x[tuple2]=r'{}\{}'.format(dir,entity)
           else:
              if '.lnk' in entity or '.LNK' in entity:
                 continue
              x[entity]=r'{}\{}'.format(dir,entity)

def getmeta(file):
   os.chdir(r'C:\database1')
   try:
      a=os.popen(r'tool.bat "{}"'.format(file))
   except:
      createtool()
      a=os.popen(r'tool.bat "{}"'.format(file))
   for  word in a:
      if 'artist' in word or 'Artist' in word:
         return word


updating=False
thislabel=None
from tool import*               
updatecount=0
def update_database():
   global box2,upbox,box4,updatecount,frame2,framen,button5,binitial,fame1,textbox
   global myButton,root,entryBox,thislabel,myButtontwo,updating
   updating=True
   if labelenter is not None:
       labelenter.destroy()
   if button5!=None:
       button5.destroy()
   if binitial!=None:
       binitial.destroy()
   entryBox.destroy()
   myButton.destroy()


   try:
      dingdong=Label(frame2,text='SCANNING...',bg='black',fg='blue')
      dingdong.pack(side=TOP)
      incpdf=Checkbutton(root,text='Index pdf docs also (Slower! )',bg='gray' ,fg='black',font='Times 10',command=flippdf)
      incpdf.pack()
      incpdf.place(relx=1,x=-80,y=2,anchor=NE)

   except:
      dingdong=Label(framen,text='SCANNING...',bg='black',fg='blue')
      dingdong.pack(side=TOP)
      incpdf=Checkbutton(root,text='Index pdf docs also (Slower! )',bg='gray' ,fg='black',font='Times 10',command=flippdf)
      incpdf.pack()
      incpdf.place(relx=1,x=-80,y=2,anchor=NE)
   if not os.path.exists(r'C:\database1'):
      os.mkdir(r'C:\database1')
   createtool()
   ll2=re.findall(r"[D-Z]+:",os.popen("wmic logicaldisk get deviceid").read(),re.MULTILINE)+re.findall(r"[A-B]+:",os.popen("wmic logicaldisk get deviceid").read(),re.MULTILINE)
   for direc in ll2:
     try:
         ll3=re.findall(r"[A-Z]+:",os.popen("wmic logicaldisk where drivetype=2 get deviceid").read(),re.MULTILINE)
         if direc in ll3:
            continue
         rcheck=os.popen("vol {}".format(direc))
         rchk=rcheck.readline()
         if 'recovery' in rchk or 'RECOVERY' in rchk:
            continue

         database(direc)
     except:
         continue
   os.chdir(r'C:\Users')
   dirinfo=os.popen('dir /B')
   for direc in dirinfo:
      str(direc)
      dirc=direc.strip()
      try:
         database(r"C:\Users\{}".format(dirc))
      except:
         continue
   database(r"C:")
   update_file_database()
   if frame2!=None:
      frame2.destroy()
   try:
      framen.destroy()
   except:
      print('no')
   dingdong.destroy()
   Label(root,text="One more step...",fg='white',bg='black').pack(side=TOP)
   thislabel=Message(root,text="The Documents are being scanned and indexed. Plz dont exit during this operation.\nWill take quite long if pdf documents are also involved.",fg='white',bg='blue',width=500)
   thislabel.pack()
   textbox=Text(root,height=25,width=150,bg='white')
   textbox.pack(fill=Y)

   thfun2.start()
   updating=False

thfun=threading.Thread(target=update_database)

def update_file_database():
   global updatecount
   if not os.path.exists(r'C:\database1'):
       os.mkdir(r'C:\database1')
   os.chdir(r'C:\database1')
   fileopen=open('DATABASEF.pkl','wb')
   pickle.dump(x,fileopen)
   fileopen.close()

var2=False


def flippdf():
    global includepdf
    if includepdf ==False:
        includepdf=True
    elif includepdf==True:
        includepdf=False
def sustainit():
   global var2
   if var2==False:
      var2=True
   elif var2==True:
      var2=False
def opennow(path):
    global root,root2,root5,root4
    os.popen('explorer.exe "{}"'.format(path))
    if var2!=True:
       root.destroy()
       try:
           root2.destroy()
       except:
           print('no')
       try:
           root5.destroy()
       except:
           print('no')
       try:
           root4.destroy()
       except:
           print('no')

frame2=None
button5=None
def upcommand():
    global frame2,root,button5,frame1,updating,frameexists,entryBox,labelenter,fram,textbox
    updating=True
    frame2=Frame(root)
    frame2.pack(side=TOP)
    frame2.config(background='gray')
    frame1.destroy()
    frameexists=False
    entryBox.destroy()
    myButton.destroy()
    frme.destroy()
    fram.destroy()
    llll.destroy()
    llll2.destroy()
    Label(frame2,text="To get started, a virtual Database will be created. ",bg='gray',fg='black',font='Mincho').pack(side=TOP)
    Label(frame2,text="Creating the database...takes a few minutes of your life... It's a one time investment.   ",bg='gray',fg='black',font='Mincho').pack(side=TOP)
    Label(frame2,text="After that, the data index is always automatically updated(Background)",bg='gray',fg='black',font='Mincho').pack(side=TOP)
    Label(frame2,text=r"Drive C: will be scanned and your custom partitions(if any) for ex. F: or whatever it is named will be scanned",font='Mincho',bg='gray',fg='black').pack(side=TOP)
    Label(frame2,text="Plz Be Patient! I will inform you when done.",bg='gray',fg='black',font='Mincho').pack(side=TOP)
    button5=Button(frame2,text="Do it!",command=thfun.start,bg='dim gray')
    button5.pack(ipadx=20)

    scroller = Scrollbar(frame2)
    scroller.pack(side=RIGHT,fill=Y)
    textbox=Text(frame2,height=25,width=150,bg='white',fg='black')
    textbox.pack(fill=Y)
    scroller.config(command=textbox.yview)
    textbox.config(yscrollcommand=scroller.set)

    if labelenter is not None:
       labelenter.destroy()

def remove(s):
   b=[',','<','.','>','/','\',','?',"'",'"',':',';','[',']','_','-','*','&','^','%','$','#','@','!','~','(','\n',')','+','=','...',"{","}","|","\\"]
   c=s.strip()
   n=len(c)
   list1=[]
   list2=[]
   for i in range(0,n):
       list1.append(c[i])
   d=0
   for i in range(0,n):
       if list1[d] in b:
          list2.append(list1[d])
          list1[d]=' '
       d=d+1
   return list1

not_destroyed=True

def splt(string):
    e=''
    s=remove(string)
    d=[]
    for alpha in s:
        e+=alpha
    b=e.split()
    for element in b:
        d.append(element.lower())
    return d

framen=None
u=0
binitial=None
dbse=None
checktest=None
from tkinter.scrolledtext import*
def initsearch():
       global box2,framen,binitial,entryBox,myButton,myButtontwo,labelenter,fram,frame1,dbse,checktest,fram,textbox,upbox,frameexists

       os.chdir(r'C:\database1')
       fileopen=open('DATABASEF.pkl','rb')
       dbse=pickle.load(fileopen)
       fileopen.close()
       for value in dbse:
           continue
       checktest=True



def search(key,frame1):
    global box2,framen,binitial,entryBox,myButton,myButtontwo,labelenter,fram,dbse,checktest,frme
    global u
    y={}
    global frameexists
    global upbox
    frameexists=True
    if searchlabel!=None:
        searchlabel.destroy()

       
    list=splt(str(key))
    n=len(list)
    k=0
    if checktest==True:
       points={}
       for word in dbse:
          for _ in list:
             if type(word) is tuple:
                if _ in word[0].lower() or key in word[0].lower():
                   k+=1
                   if (word,dbse[word]) in points:
                     points[(word,dbse[word])]+=1
                   else:
                     points[(word,dbse[word])]=1
                   continue
                if word[1] is not None and _ in word[1].lower():
                    k+=1
                    if (word,dbse[word]) in points:
                       points[(word,dbse[word])]+=1
                    else:
                       points[(word,dbse[word])]=1
             if type(word) is str:
                if _ in word.lower() or key in word.lower():
                   k+=1
                   if (word,dbse[word]) in points:
                     points[(word,dbse[word])]+=1
                   else:
                     points[(word,dbse[word])]=1


    stac=0
    sequence=0
    sequence2=0
    for key,value in sorted(points.items(),key=operator.itemgetter(1),reverse=True):
        try:
          stac+=1
          if type(key[0]) is tuple:
              initial=key[0][0]
              final=key[1]
              indx=key[0][0].lower().find(list[0].lower())
          if type(key[0]) is str:
              initial=key[0]
              final=key[1]
              indx=key[0].lower().find(list[0].lower())
          final2=os.path.dirname(final)
          
          if ".mp3" in initial or ".MP3" in initial or ".MP4" in initial or ".3gp" in initial or ".mp4" in initial or ".flv" in initial or ".mkv" in initial or ".wav" in initial or ".avi" in initial or ".jpg" in initial or ".jpeg" in initial or ".gif" in initial or ".png" in initial or '.JPG' in initial:
             button=Button(frme,underline=indx,cursor='hand2',highlightcolor='blue',text=initial,command=partial(opennow,final),anchor=W,bg="white",relief=GROOVE,activebackground="SteelBlue2",height=1,width=80)
             button.pack(side=TOP)  
             button.place(relx=1,x=-25,y=sequence2,anchor=NE)
             button2=Button(frme,cursor='hand2',text="Location",command=partial(opennow,final2),bg="SteelBlue2",relief=FLAT,activebackground="gray53",height=1)
             button2.pack()
             button2.place(relx=1,x=-9,y=sequence2,anchor=NE)
             sequence2+=26
          else:
             button=Button(frame1,underline=indx,cursor='hand2',highlightcolor='blue',text=initial,command=partial(opennow,final),anchor=W,bg="white",relief=GROOVE,activebackground="SteelBlue2",height=1,width=80)
             button.pack(side=TOP)
             button.place(relx=1,x=-89,y=sequence,anchor=NE)
             button2=Button(frame1,cursor='hand2',text="Open file location",command=partial(opennow,final2),bg="SteelBlue2",relief=FLAT,activebackground="gray53",height=1)
             button2.pack()
             button2.place(relx=1,x=-9,y=sequence,anchor=NE)
             if os.path.isdir(final):
                some=Label(fram,text="FILE FOLDER",bg="burlywood2",anchor=E,width=28,height=1)
                some.pack()
                some.place(relx=1,x=-60,y=sequence,anchor=NE)

             elif ".pdf" in initial:
                some=Label(fram,text="pdf Document",anchor=E,bg="firebrick1",width=28,height=1)
                some.pack()
                some.place(relx=1,x=-60,y=sequence,anchor=NE)
             elif ".exe" in initial:
                some=Label(fram,text="Executable",anchor=E,bg="gray53",width=28,height=1)
                some.pack()
                some.place(relx=1,x=-60,y=sequence,anchor=NE)
             elif ".txt" in initial:
                some=Label(fram,text="Text Document",anchor=E,bg="antiquewhite1",width=28,height=1)
                some.pack()
                some.place(relx=1,x=-60,y=sequence,anchor=NE)
             elif ".html" in initial or ".htm" in initial:
                some=Label(fram,text="HTML Document",anchor=E,bg="DarkSeaGreen1",width=28,height=1)
                some.pack()
                some.place(relx=1,x=-60,y=sequence,anchor=NE)
             else:
                some=Label(fram,text="Document",anchor=E,bg="floral white",width=28,height=1)
                some.pack()
                some.place(relx=1,x=-60,y=sequence,anchor=NE)
             sequence+=26
          if stac==30:
              break
        except:
          continue
    global lkill,label2

    if k==0:
           label2=Label(root,text="No results",bg="gainsboro",fg="gray",width=120)
           label2.pack()
           label2.place(height=10,width=120,x=2,y=2)
           clearf(frame1);clearf(fram);clearf(frme)
           lkill=True
    if k!=0 and lkill==True:
           label2.destroy()
           yo='No'



def clearf(frame):
    a=frame.winfo_children()
    for widget in a:
        widget.destroy()
txt=None
itxt=None
def buttonPushed():
   global entryBox,txt,itxt,updating
   global root
   global frame1
   global frameexists
   if txt!=None:
      itxt=txt
   txt = entryBox.get()
   if txt=="":
      clearf(frame1)
      clearf(frme)
      clearf(fram)
   if txt!='' and itxt!=txt and txt!=' ' and not updating:
        search(txt,frame1)

def buttonpushed2():
     global entryBox,txt,itxt,thinitbut
     global root
     global frame1
     global frameexists
     if txt!=None:
        itxt=txt
     txt = entryBox.get()
     if txt!='' and itxt!=txt and txt!=' ':
          if frameexists:
             info=frame1.winfo_children()
             for widget in info:
                widget.destroy()
             clearf(fram)
          search(txt,frame1)

def bp2(event):
   threading.Thread(target=buttonpushed2).start()

from time import sleep
def initbut():
        global txt,mailinitiated,updating,grint,labelenter,itxt
     
        while True:
           try:
              if mailinitiated or updating or grint:
                  labelenter=Label(root,text="Press Enter To Search   ",fg='gray',bg="gainsboro",height=3,width=30)
                  labelenter.pack()
                  labelenter.place(height=10,width=120,y=2,x=3)
                  break
              buttonPushed()
              sleep(0.1)
           except:
              continue





def createTextBox(parent):
   global entryBox
   entryBox = Entry(parent,background="gainsboro",width=60,relief=SUNKEN,font='Mincho')
   entryBox.pack(ipady=4,side=TOP)
   entryBox.place(height=40,width=400,x=0,y=0)
   entryBox.focus_set()
   entryBox.bind('<Return>',bp2)


u=0
v=0
x={}
def checksum(dir):   #a function that returns the checksum of a given file of any format
    a=os.popen('certutil -hashfile "{}"'.format(dir))
    for word in a:
       b=a.readline()
       return b

completed=False
textbox2=None
def check_dir(dire):
    global v,root1,textbox2
    global u,completed

    b=os.listdir(os.chdir(dire))   #b is a list
    for entity in b:   #entity is the current folder or file being processed
        if os.path.isdir("{}\{}".format(dire,entity)):
            if len(os.listdir("{}\{}".format(dire,entity)))==0:
                os.rmdir("{}\{}".format(dire,entity))
                v+=1
                Label(root1,text="The folder {}\{} has been deleted because it was empty".format(dire,entity)).pack(side=TOP)
                textbox2.insert(END,"The folder {}\{} has been deleted because it was empty\n".format(dire,entity))
                textbox2.see(END)
                continue
            check_dir("{}\{}".format(dire,entity))


        else:
            c=checksum("{}\{}".format(dire,entity))
            if c in x:
                x[c]+=1
                os.remove("{}\{}".format(dire,entity))
                u+=1
                #print("the file {}\{} has been deleted".format(dir,entity))
                textbox2.insert(END,"the file {}\{} has been deleted\n".format(dire,entity))
                textbox2.see(END)
            else:
                x[c]=1
    completed=True
    x.clear()
l1=None
l2=None
proceed=None
dupinitiated=False
def init_dup():
   global root1,b,proceed,dupinitiated
   root1=Tk()
   root1.title("Remove duplicate files")
   dupinitiated=True
   global box6


   Message(root1,text="WARNING: Please note that any empty folders and empty files will also be deleted. \nThe app may hang during this operation",bg='antique white',width=500).pack(side=TOP)
   proceed=Label(root1,text="Do you want to proceed? [Y/N]")
   proceed.pack(side=TOP)
   box6=Entry(root1)
   box6.pack(side=TOP)
   box6.focus_set()
   b=Button(root1,text='Go On',bg='gray53',activebackground='peach puff')
   b.pack()
   b.bind('<Button-1>',get_response)
   root1.mainloop()
but=None


def get_response2():
    global g,box6,root1,not_destroyed,but,l1,l2,root,textbox2

    h=box6.get()
    x={}
    textbox2=Text(root1)
    textbox2.pack()
    try:
        check_dir(h)   #calling the function
        box6.destroy()
        but.destroy()
        l1.destroy()
        l2.destroy()
        if not_destroyed:
           if u!=0 or v!=0:  #checking if any file was actually deleted
              Message(root1,text='Total {} duplicate file(s) and {} empty folder(s) was/were deleted.\nPlease Exit the whole application.\nPlease do not use this feature more than once in one session.'.format(u,v),width=300).pack(side=TOP)
           elif completed:
              Message(root1,text="Scan complete.\nNo duplicate files were found.\nNow please Exit the whole application.\nPlease do not use this feature more than once in one session.",width=300).pack(side=TOP)
        exitbu=Button(root1,text='EXIT',bg='gray53',activebackground="peachpuff")
        exitbu.pack(side=TOP)
        exitbu.focus_set()
        exitbu.bind('<Return>',dest)
        exitbu.bind('<Button-1>',dest)

    except:
        Label(root1,text="Path Not found.Operation failed.Please restart the application").pack(side=TOP)

def dest(event):
      os._exit(0)


def get_response(event):

    global box6,g,b,root1,u,v,but,proceed,grint,gtxt,l1,l2
    global not_destroyed
    u=0
    v=0
    grint=True
    box6.focus_set()
    g=box6.get()
    box6.delete(0,END)
    proceed.destroy()
    if g=='Y' or g=='y':
       l1=Label(root1,text="Enter the drive or directory(full path) in which")
       l2=Label(root1,text="you want to remove duplicated files.For Ex. 'F:'")
       l1.pack()
       l2.pack()

       b.destroy()
       b=None

       but=Button(root1,text='DO IT',bg='gray53',activebackground="peach puff",command=get_response2)
       but.pack(side=TOP)


    elif g=='N' or g=='n':
        root1.destroy()
    else:
        Label(root1,text="Unrecognized response. Please try again")
        root1.destroy()
        init_dup()
        not_destroyed=False


def removeit(s):
   b=[',','<','.','>','/','\',','?',"'",'"',':',';','[',']','_','-','*','&','^','%','$','#','@','!','~','(','\n',')','+','=','...',"{","}","|","\\"]
   n=len(s)
   list1=[]
   for i in range(0,n):
       list1.append(s[i])
   c=0
   for i in range(0,n):
       if list1[c] in b:
          del list1[c]
          continue
       c=c+1
   return list1

yugbox=None
yugbox2=None
yugframe=None
def yughere(event):
   global yugbox,yugbox2,yugframe
   iamtrue=True
   yug1=str(yugbox.get())
   yug2=str(yugbox2.get())
   try:
      a=open(yug1)
      d=open(yug2,'x')
   except:
      Label(yugframe,text='Error: Either the path is incorrect or the second file already exists.').pack(side=TOP)
      iamtrue=False
   if iamtrue:
     try:
         for word in a:
          c=removeit(word)
          for l in c:
              d.write(l)

          d.write("\n")
     except:
         Label(yugframe,text='Something went wrong.Please Try again')
   if iamtrue:
      clearf(yugframe)
      Label(yugframe,text="Ok, I'm done! Let's now meet the unhappy file (:P) ",height=5,bg="SteelBlue2").pack()



def inityughere():
    global yugbox,yugbox2,yugframe
    root2=Tk()
    root2.config(background='RoyalBlue4')
    yugframe=Frame(root2,bg='RoyalBlue4')
    yugframe.pack()
    Label(yugframe,text="Enter the full path to the text file(.txt extension)",bg='RoyalBlue4',fg='white').pack(side=TOP)
    yugbox=Entry(yugframe,bg='lemon chiffon')
    yugbox.pack(side=TOP)
    Label(yugframe,text="Now enter the full path (including file name which will be created automatically) to store the 'unhappy' text",fg='white',bg='RoyalBlue4').pack(side=TOP)
    yugbox.focus_set()
    yugbox2=Entry(yugframe,bg='lemon chiffon')
    yugbox2.pack(side=TOP)
    yugbox.bind('<Return>',focusSET)
    yugbox2.bind('<Return>',yughere)
    yugbutton=Button(yugframe,text="Do This Crap!",bg='gray53')
    yugbutton.pack(side=TOP)
    yugbutton.bind('<Button-1>',yughere)

def focusSET(event):
   global yugbox2
   yugbox2.focus_set()


def removenew(s):
   thisb=[',','<','.','>','/','?',"'",'"',':',';','[',']','_','-','*','&','^','%','$','#','@','!','~','(',')','+','=','...',"{","}","|","\\"]
   n=len(s)
   thislist1=[]
   for i in range(0,n):
       thislist1.append(s[i])
   c=0
   for i in range(0,n):
       if thislist1[c] in thisb:
          del thislist1[c]

          continue
       c=c+1

   string=''

   for ele in thislist1:
      string+=ele
   thislist2=string.split()
   return thislist2


def scan2():
   global thislabel,updating,textbox
   updating=True
   stopwordlist=['None','is','an','the','are','a','of','and','to','for','in','it','',' ']
   pdfdict={}
   thisfile1=open(r'C:\database1\DATABASEF.pkl','rb')
   dbse=pickle.load(thisfile1)
   xdct={}

   for entity in dbse:
     sleep(0.005)

     if '.docx' in entity:

        try:
           try:
              textbox.insert(END,"{}\n".format(entity))
              textbox.see(END)
           except:
              print('No')
           a=docx.Document(dbse[entity])
           for line in a.paragraphs:
              b=removenew(line.text)
              for word in b:
                 if word=='' or word==' ':
                    continue
                 elif '\u2019' in word:
                    wordn=word.replace('\u2019','')
                 elif '\u2026' in word:
                    wordn=word.replace('\u2026','')
                 else:
                    wordn=word
                 if wordn in xdct:
                    if dbse[entity] in xdct[wordn]:
                       continue
                    xdct[wordn].append(dbse[entity])
                    continue
                 listm=[]
                 listm.append(dbse[entity])
                 xdct[wordn]=listm
                 
        except:
           continue

     if '.txt' in entity or '.csv' in entity or '.srt' in entity:
        try:
           try:
              textbox.insert(END,"{}\n".format(entity))
              textbox.see(END)
           except:
              print('No')
           am=open(dbse[entity])

           for word in am:
               list2=removenew(word)
               for e in list2:
                  d=e.replace('\x00','')
                  if d=='' or d==' ' or '1' in d or '2' in d or '3' in d or '4' in d or '5' in d or '6' in d or '7' in d or '8' in d or '9' in d or '0' in d or d in stopwordlist:
                     continue
                  if d in xdct:
                     if dbse[entity] in xdct[d]:
                        continue
                     xdct[d].append(dbse[entity])
                     continue
                  listn=[]
                  listn.append(dbse[entity])
                  xdct[d]=listn

        except:
           continue

     if '.pptx' in entity:

        try:
            try:
              textbox.insert(END,"{}\n".format(entity))
              textbox.see(END)
            except:
              print('No')
            prs = Presentation(dbse[entity])
            for slide in prs.slides:
               for shape in slide.shapes:
                  if not shape.has_text_frame:
                     continue
                  for paragraph in shape.text_frame.paragraphs:
                     a=removenew(paragraph.text)
                     for ptword in a:
                        if ptword in stopwordlist:
                           continue
                        if ptword in xdct:
                            if dbse[entity] in xdct[ptword]:
                               continue
                            xdct[ptword].append(dbse[entity])
                            continue
                        listn=[]
                        listn.append(dbse[entity])
                        xdct[ptword]=listn
        except:
           continue


     if '.pdf' in entity and includepdf:
         try:
            pdfdict[entity]=dbse[entity]
         except:
            print("no")
         



     if '.xlsx' in entity:
        try:
           try:
              textbox.insert(END,"{}\n".format(entity))
              textbox.see(END)
           except:
              print('No')
           wb=openpyxl.load_workbook(dbse[entity])
           getnm=wb.get_sheet_names()
           for sheetn in getnm:
              sheet=wb.get_sheet_by_name(sheetn)

              hr=sheet.max_row
              hc=sheet.max_column
              
              for i in range(1,hr):
                 sleep(0.005)
                 for j in range(1,hc):
                    clvalue2=sheet.cell(row=i,column=j).value
                    clvalue="{}".format(clvalue2)
                    if type(clvalue) is not int:
                       if clvalue in stopwordlist or clvalue is None or '1' in clvalue or '2' in clvalue or '3' in clvalue or '4' in clvalue or '5' in clvalue or '6' in clvalue or '7' in clvalue or '8' in clvalue or '9' in clvalue or '0' in clvalue:
                           continue
                    else:
                       continue
                    if type(clvalue) is str:
                       clvalue2=removenew(clvalue)
                       for ant in clvalue2:
                           if ant in xdct:
                               xdct[ant].append(dbse[entity])
                               continue
                           listn=[]
                           listn.append(dbse[entity])
                           xdct[ant]=listn
                        
        except:
            continue

   thisfile2=open('C:\database1\DATABASEW.pkl','wb')
   pickle.dump(xdct,thisfile2)
   thisfile2.close()
   for entity in pdfdict:
      if includepdf:
             try:
                textbox.insert(END,"{}\n".format(entity))
                textbox.see(END)
             except:
                print('No')
             fpn = open(dbse[entity], 'rb')
             parser = PDFParser(fpn)
             docm = PDFDocument()
             parser.set_document(docm)
             try:
                docm.set_parser(parser)
                docm.initialize('')
                rsrcmgr = PDFResourceManager()
                laparams = LAParams()
                device = PDFPageAggregator(rsrcmgr, laparams=laparams)
                interpreter = PDFPageInterpreter(rsrcmgr, device)
                for page in docm.get_pages():
                    interpreter.process_page(page)
                    layout = device.get_result()
                    for lt_obj in layout:
                        if isinstance(lt_obj, LTTextBox) or isinstance(lt_obj, LTTextLine):
                            try:
                              for wrd in removenew(lt_obj.get_text()):
                                 sleep(0.005)
                                 if '1' in wrd or '2' in wrd or '3' in wrd or '4' in wrd or '5' in wrd or '6' in wrd or '7' in wrd or '8' in wrd or '9' in wrd or '0' in wrd:
                                    continue
                                 if type(wrd) is str:

                                     if wrd in stopwordlist:
                                         continue
                                     if wrd in xdct:
                                        xdct[wrd].append(dbse[entity])
                                        continue
                                     listn=[]
                                     listn.append(dbse[entity])
                                     xdct[wrd]=listn


                            except:
                              continue
             except:
                 continue
      else:
             continue

   thisfile2=open('C:\database1\DATABASEW.pkl','wb')
   pickle.dump(xdct,thisfile2)
   thisfile2.close()
   thisfile1.close()
   thislabel2=Label(text="Great! Now search inside of all kinds of text docs as well. Please restart application",fg='white',bg='blue')
   thislabel2.pack()
   thislabel.destroy()



thfun2=threading.Thread(target=scan2)

altsearch=False
suplist=None
def search2(event):
    global frame1,frameexists,thisdict2,altsearch,suplist,fram,frme,llll2
    dict3={}
    altsearch=True
    if frameexists:
       clearf(frame1)
    clearf(frme)
    clearf(fram)
    llll2.destroy()
    txt = entryBox.get()
    thislist=[]
    if txt !='':
       for element in txt.split():
           elelist=capspermutation(element)
           for _ in elelist:
               thislist.append(_)

       for words in thislist:
           try:
              if words in thisdict2 and words!='':
                  if words in dict3:
                      continue
                  dict3[words]=thisdict2[words]
           except:
               continue

       n=len(thislist)
       e=n
       suplist={}
       count=0
       boogle=0
       for word in dict3:
           if len(dict3[word])<100:
              for ent in dict3[word]:
                  if ent in suplist:
                      suplist[ent]+=1
                      continue
                  suplist[ent]=1

       for key,value in sorted(suplist.items(),key=operator.itemgetter(1),reverse=True):
           if count<22:
               butt=Button(frame1,cursor='hand2',text=os.path.basename(key),command=partial(opennow,key),bg='white',fg='black',width=83,activebackground='gray53',relief=GROOVE,font='Mincho 10')
               butt.pack(side=TOP)
               count+=1
           if count>=22:
               butt=Button(frme,cursor='hand2',text=os.path.basename(key),command=partial(opennow,key),bg='white',fg='black',width=83,activebackground='gray53',relief=GROOVE,font='Mincho 10')
               butt.pack(side=TOP)
               if boogle==22:
                   count=0
                   break
               boogle+=1


def seeall():
    global suplist,root5
    if len(suplist)!=0:
       root5=Tk()
       root5.title('See All')
       root5.config(background='peach puff')
    count=0
    for key,value in sorted(suplist.items(),key=operator.itemgetter(1),reverse=True):
           Button(root5,text=os.path.basename(key),command=partial(opennow,key),bg='gray',fg='white',width=100,activebackground='gray53',relief=GROOVE,font='Mincho 10').pack(side=TOP,ipadx=50)
           count+=1
           if count==29:
               count=0
               break




thisvar=None
thisdict2=None
def initiates2():
    global thisdict2,thisvar,frameexists
    try:
       os.chdir(r"C:\database1")
       thisvar=open(r"C:\database1\DATABASEW.pkl",'rb')
       thisdict2=pickle.load(thisvar)
    except:
       if frameexists:
          upcommand()



import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email.mime.text import MIMEText
from email.utils import COMMASPACE, formatdate
from email import encoders

def send_mail( emailprovider,send_from, send_to, subject, text, files=[], port=587, username='', password='', isTls=True):
    msg = MIMEMultipart()
    msg['From'] = str(send_from)
    msg['To'] = COMMASPACE.join(str(send_to))
    msg['Date'] = formatdate(localtime = True)
    msg['Subject'] = str(subject)

    msg.attach( MIMEText(text) )

    for f in files:
        part = MIMEBase('application', "octet-stream")
        part.set_payload( open(f,"rb").read() )
        encoders.encode_base64(part)
        part.add_header('Content-Disposition', 'attachment; filename="{0}"'.format(os.path.basename(f)))
        msg.attach(part)

    smtp = smtplib.SMTP('smtp.{}'.format(emailprovider), port)
    if isTls: smtp.starttls()
    smtp.login(username,password)
    smtp.sendmail(send_from, send_to, msg.as_string())
    smtp.quit()






root4=None
mframe=None
mailinitiated=False
var=None
def init_mail():
    global root4,mframe,mailinitiated,var
    mailinitiated=True
    root4=Tk()
    root4.title('Quick Email Facility')
    root4.config(background='SteelBlue1')
    root4.resizable(width=False,height=False)
    mframe=Frame(root4)
    mframe.pack()
    mframe.config(bg='SteelBlue2')
    var=0
    Message(mframe,text="To use this quick email facility, please make sure that you have \n enabled the 'Allow less secure Apps' option in gmail account settings. \n\nMy account-> Sign in And security-> Allow less Secure apps \n\nYour mailing address and password will never be stored.\nIt is totally safe.",bg='SteelBlue2',fg='black',font='Times',width=600).pack()
    chk=Checkbutton(mframe,text='Email With Attachment',command=variable,bg='SteelBlue2')
    chk.pack()
    Button(mframe,text='Done !',command=go_mail,bg='gray53',activebackground='peach puff',relief=GROOVE).pack(ipadx=20)
    root4.mainloop()

def variable():
    global var
    if var==0:
        var=1

    else:
        var=0
ethread=None
e1=None;e2=None;e3=None;e4=None;e5=None;e6=None;eframe=None;e7=None
def go_mail():
    global mframe,root4,e1,e2,e3,e4,e5,e6,e7,eframe,var,ethread
    mframe.destroy()
    eframe=Frame(root4)
    eframe.pack()
    eframe.config(background='SteelBlue1')
    e1=Entry(eframe,background='peach puff',font='Times')
    l1=Label(eframe,text='Email:')
    l1.pack(ipadx=60,ipady=4)
    e1.pack(ipadx=60,ipady=4)
    e2=Entry(eframe,background='peach puff',font='Times',show='*')
    l1=Label(eframe,text='Password')
    l1.pack(ipadx=60,ipady=4)
    e2.pack(ipadx=60,ipady=4)
    e3=Entry(eframe,background='peach puff',font='Times')
    l1=Label(eframe,text='Reciever(s):')
    l1.pack(ipadx=60,ipady=4)
    e3.pack(ipadx=60,ipady=4)
    e4=Entry(eframe,background='peach puff',font='Times')
    l1=Label(eframe,text='Subject:')
    l1.pack(ipadx=60,ipady=4)
    e4.pack(ipadx=60,ipady=4)
    e7=Entry(eframe,background='peach puff',font='Times')
    l1=Label(eframe,text='Your Email provider(For example-  gmail.com)')
    l1.pack(ipadx=60,ipady=4)
    e7.pack(ipadx=60,ipady=4)
    e5=Text(eframe,background='peach puff',font='Times',height=10,width=60)
    l1=Label(eframe,text='Body:')
    l1.pack(ipadx=60,ipady=4)
    e5.pack(ipadx=60,ipady=4)
   
    if var:
       e6=Entry(eframe,background='peach puff',font='Times')
       l1=Label(eframe,text='Attachment file(Full Path To The File)')
       l1.pack(ipadx=60,ipady=4)
       e6.pack(ipadx=60,ipady=4)
    ebutton=Button(eframe,text='send',bg='lightgreen',command=makesit,activebackground='peach puff',relief=GROOVE,font='Times')
    ebutton.pack(ipadx=20)




elabel3=None
elabel2=None
elabel3made=False;elabel2made=False
def sendit():
    global eframe,e1,e2,e3,e4,e5,e6,e7,elabel3made,elabel2made,elabel3,elabel2,var,ethread
    ethread='started'
    efrom=e1.get()
    epassword=e2.get()
    eto=e3.get()
    esubject=e4.get()
    etext=e5.get(index1=1.0,index2=END)
    emailprovider=e7.get()
    if var:
       eattach=e6.get()
    if elabel3!=None:
        try:
           elabel3.destroy()
        except:
           elabel3=None
    if elabel2!=None:
       try:
         elabel2.destroy()
       except:
         elabel2=None
    elabel=Label(eframe,text='Sending mail. Please Wait.')
    elabel.pack()
    try:
       if var:
          send_mail(emailprovider,efrom, eto,esubject, etext,files=[eattach],username=efrom, password=epassword)
       else:
          send_mail2(emailprovider,efrom,eto,efrom,epassword,esubject,etext)
       elabel.destroy()
       elabel2=Label(eframe,text='Sent Successfully!',font='Times')
       elabel2.pack()
       elabel2made=True
    except:
       elabel.destroy()
       elabel3=Label(eframe,text='Sending Failed')
       elabel3.pack()
       elabel3made=True



def send_mail2(emailprovider,sender,reciever,username,password,subject,text):
   message="subject: {} \n{}".format(subject,text)
   a=smtplib.SMTP('smtp.{}'.format(emailprovider),25)
   a.ehlo()
   a.starttls()
   a.login(username,password)
   a.sendmail(sender,reciever,message)


def makesit():
   sit=threading.Thread(target=sendit)
   sit.start()



searchlabel=None
thisdict2=None
muButton=None
myButtontwo=None
thinitbut=None

def hlp():
   root5=Tk()
   Label(root5,text='Welcome to seachBASE 1.5 , a simple, easy to setup and use, Desktop search software.').place(x=40,y=0,anchor=NW)
   Label(root5,text='Following are the specifications and usage directions:').place(x=40,y=30,anchor=NW)
   Label(root5,text=r'-To begin, it indexes the files,folders and the documents inside the drives.').place(x=70,y=60,anchor=NW)
   Label(root5,text='-After that, the indexing becomes automatic and a background process.').place(x=70,y=80,anchor=NW)
   Label(root5,text='-Once you start entering the keyword, the appropriate results start appearing').place(x=70,y=100,anchor=NW)
   Label(root5,text='-The "Search inside Documents" button puts forward all those docmuents (including pdf) that contain the entered keyword(s); with a predefined priority order.').place(x=70,y=120,anchor=NW)
   Label(root5,text='-The Quick Email facility requires you to enter your email and password and other details').place(x=70,y=140,anchor=NW)
   Label(root5,text='-In any case, any of the details(like the password or the email) will NEVER be stored. It is totally safe.').place(x=70,y=160,anchor=NW)
   Label(root5,text='-The "remove duplicate files" removes the duplicate files from your desired folder or drive, WITHOUT notifying you about the names of the duplicate files deleted').place(x=70,y=180,anchor=NW)
   Label(root5,text='-However this is safe but still this utility should NOT be used twice in the same session.').place(x=70,y=200,anchor=NW)

SIS=False
def main():
   global root,frame1,myButton,thisdict2,myButtontwo,thinitbut,frameexists,searchlabel,fram,frme,llll,llll2,SIS
   supremedir=os.getcwd()
   initiates2()
   try:
      initsearch()
   except:
      SIS=True
   root = Tk()
   root.title("searchBASE")
   root.configure(background="gray53")
   widthpixels=1365
   heightpixels=710
   root.geometry('{}x{}'.format(widthpixels, heightpixels))
   root.resizable(height=False,width=False)
   os.chdir(supremedir)
   try:
      root.iconbitmap('search2.ico')
   except:
      try:
         root.iconbitmap(r'C:\Program Files\searchBASE\search2.ico')#for the icon
      except:
         print("no")
   myButton = Button(root,cursor='hand2', text="Search inside Documents",bg="SteelBlue1",relief=GROOVE,activebackground="white",width=19)
   myButton.pack(side=TOP)
   sust=Checkbutton(root,cursor='hand2',text='Sustain after clicking a result',bg='gray53' ,fg='black',font='Times 10',command=sustainit)
   sust.pack()
   sust.place(relx=1,x=-80,y=2,anchor=NE)

   createTextBox(root)
   llll=Label(root,text='Documents',width=108,anchor=W,bg='white',font='Mincho 10')
   llll.pack(ipadx=10)
   llll.place(x=0,y=50)
   llll2=Label(root,text='Media files',width=108,anchor=W,bg='light gray',font='Mincho 10')
   llll2.pack(ipadx=10)
   llll2.place(x=760,y=50)
   frme=Frame(root,height=700,width=600)
   frme.pack()
   frme.place(x=762,y=70)
   frme.config(background='light gray')
   frme.pack_propagate(False)

   fram=Frame(root,height=700,width=160)
   fram.pack(side=LEFT)
   fram.place(x=0,y=70)
   fram.config(background='white')
   fram.pack_propagate(False)

   frame1=Frame(root,height=700,width=660)
   frame1.pack()
   frame1.place(x=100,y=70)
   frame1.config(background='white')
   frame1.pack_propagate(False)



   searchlabel=Message(frame1,bg='white',fg='gray53',text="`Search just got smarter \n`Type an artist's name and their creations will be shown.\n`Common Stopwords are not considered for faster search results while searching documents\n`Press 'search inside documents' to search for the keyword inside your documents",width=1000,font='Times 10')
   searchlabel.pack(side=TOP)

   mymenu=Menu(root)
   frameexists=True

   utils=Menu(mymenu,tearoff=0)
   thinitbut=threading.Thread(target=initbut)
   thinitbut.start()
   utils.add_command(label='Quick email',activebackground='SteelBlue1',background='peach puff',command=init_mail)
   utils.add_separator()
   utils.add_command(label='Remove Duplicate Files',activebackground='SteelBlue1',background='peach puff',command=init_dup)
   utils.add_separator()
   utils.add_command(label='Text File Punctuation Remover',activebackground='SteelBlue1',background='peach puff',command=inityughere)
   utils.add_separator()
   utils.add_separator()
   utils.add_command(label='Help',activebackground='SteelBlue1',background='blue',command=hlp,foreground='white')
   mymenu.add_cascade(label='OPTIONS',menu=utils)
   root.config(menu=mymenu)
   myButton.bind('<Button-1>',search2)
   if SIS:
       upcommand()
   root.mainloop()
main()


try:
   tlist=os.popen('tasklist /FI "IMAGENAME eq pythonw.exe"')
   mypid2=os.getpid()
   mypid=str(mypid2)
   i=0
   alist=[]
   for line in tlist:
       alist.append(line)
   
   tlist=os.popen('tasklist /FI "IMAGENAME eq searchBASE.exe"')
   if len(alist)>4:
      for line in tlist:
          i+=1
          if i>3:
             b=line.split()
             if b[1]==mypid:
                os.popen("taskkill /F /PID {}".format(b[1]))

except:
    print('No')



from autoupdateDATA import* #to call the auto updation of data in the background



















