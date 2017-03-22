#The automatic updation of data in background takes place in this module

import os
from tool import*
from pptx import Presentation
import openpyxl
import pickle
import docx
import re
from time import*
import threading
from pdfparser import PDFParser, PDFDocument
from pdfinterp import PDFResourceManager, PDFPageInterpreter
from converter import PDFPageAggregator
from layout import LAParams, LTTextBox, LTTextLine


extlist=('.mp3','.ace','.mp4','.avi','.flv','.mkv','.jpg','.jpeg','.htm','.html','.pdf','.png','.exe','.txt','.docx','.xml','.gif','.odp','.pptx','.zip','.py','.cpp','.c','.rar','.xlsx')
x={}
def database(dir):

    b=os.listdir(dir)
    for entity in b:
        sleep(0.01)
        if entity=="$Recycle.Bin" or entity=="$RECYCLE.BIN" or entity=="System Volume Information" or "appdata" in entity.lower():
           continue
        if os.path.isdir(r'{}\{}'.format(dir,entity)):
           x[entity]=r'{}\{}'.format(dir,entity)

           try:
             database(r'{}\{}'.format(dir,entity))
           except:
             continue
        else:
           if '.mp3' in entity:
              if '.lnk' in entity or '.LNK' in entity:
                 continue
              meta=getmeta(r'{}\{}'.format(dir,entity))
              tuple2=(entity,meta)
              x[tuple2]=r'{}\{}'.format(dir,entity)
           else:
              if '.lnk' in entity or '.LNK' in entity:
                 continue
              x[entity]=r'{}\{}'.format(dir,entity)

def getmeta(file):
   os.chdir(r'C:\database1')
   try:
      a=os.popen(r'tool.bat "{}"'.format(file))
   except:
      createtool()
      a=os.popen(r'tool.bat "{}"'.format(file))
   for  word in a:
      if 'artist' in word or 'Artist' in word:
         return word


def update_database():

   if not os.path.exists(r'C:\database1'):
      os.mkdir(r'C:\database1')
   createtool()
   ll2=re.findall(r"[D-Z]+:",os.popen("wmic logicaldisk get deviceid").read(),re.MULTILINE)+re.findall(r"[A-B]+:",os.popen("wmic logicaldisk get deviceid").read(),re.MULTILINE)
   for direc in ll2:
      ll3=re.findall(r"[A-Z]+:",os.popen("wmic logicaldisk where drivetype=2 get deviceid").read(),re.MULTILINE)
      if direc in ll3:
         continue
      try:
         rcheck=os.popen("vol {}".format(direc))
         rchk=rcheck.readline()
         if 'recovery' in rchk or 'RECOVERY' in rchk:
            continue
         try:
            database(direc)

         except:
            continue
      except:
         continue
   os.chdir(r'C:\Users')
   dirinfo=os.popen('dir /B')
   for direc in dirinfo:
      str(direc)
      dirc=direc.strip()
      try:
         database(r"C:\Users\{}".format(dirc))
      except:
         continue

   update_file_database()
   scan2()



def update_file_database():

   if not os.path.exists(r'C:\database1'):
       os.mkdir(r'C:\database1')
   os.chdir(r'C:\database1')
   fileopen=open('DATABASEF.pkl','wb')
   pickle.dump(x,fileopen)
   fileopen.close()


def removenew(s):
   thisb=[',','<','.','>','/','?',"'",'"',':',';','[',']','_','-','*','&','^','%','$','#','@','!','~','(',')','+','=','...',"{","}","|","\\"]
   n=len(s)
   thislist1=[]
   for i in range(0,n):
       thislist1.append(s[i])
   c=0
   for i in range(0,n):
       if thislist1[c] in thisb:
          del thislist1[c]

          continue
       c=c+1

   string=''

   for ele in thislist1:
      string+=ele
   thislist2=string.split()
   return thislist2



def scan2():
   stopwordlist=['None','is','an','the','are','a','of','and','to','for','in','it','',' ']

   thisfile1=open(r'C:\database1\DATABASEF.pkl','rb')
   dbse=pickle.load(thisfile1)
   xdct={}

   for entity in dbse:
     sleep(0.005)

     if '.docx' in entity:
        sleep(0.005)
        try:
           a=docx.Document(dbse[entity])
           for line in a.paragraphs:
              sleep(0.005)
              b=removenew(line.text)
              for word in b:
                 if word=='' or word==' ':
                    continue
                 elif '\u2019' in word:
                    wordn=word.replace('\u2019','')
                 elif '\u2026' in word:
                    wordn=word.replace('\u2026','')
                 else:
                    wordn=word
                 if wordn in xdct:
                    if dbse[entity] in xdct[wordn]:
                       continue
                    xdct[wordn].append(dbse[entity])
                    continue
                 listm=[]
                 listm.append(dbse[entity])
                 xdct[wordn]=listm

        except:
           continue

     if '.txt' in entity or '.cpp' in entity or '.csv' in entity :
        sleep(0.005)
        try:
           am=open(dbse[entity])

           for word in am:
               list2=removenew(word)
               for e in list2:
                  d=e.replace('\x00','')
                  if d=='' or d==' ' or '1' in d or '2' in d or '3' in d or '4' in d or '5' in d or '6' in d or '7' in d or '8' in d or '9' in d or '0' in d or d in stopwordlist:
                     continue
                  if d in xdct:
                     if dbse[entity] in xdct[d]:
                        continue
                     xdct[d].append(dbse[entity])
                     continue
                  listn=[]
                  listn.append(dbse[entity])
                  xdct[d]=listn

        except:
           continue

     if '.pptx' in entity:
        try:
            prs = Presentation(dbse[entity])
            for slide in prs.slides:
               sleep(0.005)
               for shape in slide.shapes:
                  if not shape.has_text_frame:
                     continue
                  for paragraph in shape.text_frame.paragraphs:
                     a=removenew(paragraph.text)
                     for ptword in a:
                        if ptword in stopwordlist:
                           continue
                        if ptword in xdct:
                            if dbse[entity] in xdct[ptword]:
                               continue
                            xdct[ptword].append(dbse[entity])
                            continue
                        listn=[]
                        listn.append(dbse[entity])
                        xdct[ptword]=listn
        except:
           continue



     if '.pdf' in entity:
         try:
             sleep(0.005)
             continue
             
             fpn = open(dbse[entity], 'rb')
             parser = PDFParser(fpn)
             docm = PDFDocument()
             parser.set_document(docm)
             try:
                docm.set_parser(parser)
                docm.initialize('')
                rsrcmgr = PDFResourceManager()
                laparams = LAParams()
                device = PDFPageAggregator(rsrcmgr, laparams=laparams)
                interpreter = PDFPageInterpreter(rsrcmgr, device)
                # Process each page contained in the document.
                for page in docm.get_pages():
                    sleep(0.005)
                    interpreter.process_page(page)
                    layout = device.get_result()
                    for lt_obj in layout:
                        if isinstance(lt_obj, LTTextBox) or isinstance(lt_obj, LTTextLine):
                            try:
                              for wrd in removenew(lt_obj.get_text()):
                                 sleep(0.005)
                                 if '1' in wrd or '2' in wrd or '3' in wrd or '4' in wrd or '5' in wrd or '6' in wrd or '7' in wrd or '8' in wrd or '9' in wrd or '0' in wrd:
                                    continue
                                 if type(wrd) is str:
                                     if wrd in stopwordlist:
                                         continue
                                     if wrd in xdct:
                                        xdct[wrd].append(dbse[entity])
                                        continue
                                     listn=[]
                                     listn.append(dbse[entity])
                                     xdct[wrd]=listn


                            except:
                              continue
             except:
                 continue
         except:
             continue
     if '.xlsx' in entity:
        if True:
           sleep(0.005)
           wb=openpyxl.load_workbook(dbse[entity])
           getnm=wb.get_sheet_names()
           for sheetn in getnm:
              sleep(0.005)
              sheet=wb.get_sheet_by_name(sheetn)
              hr=sheet.max_row
              hc=sheet.max_column

              for i in range(1,hr):
                 sleep(0.005)
                 for j in range(1,hc):
                    eminem=False
                    clvalue2=sheet.cell(row=i,column=j).value
                    if type(clvalue2)==int or type(clvalue2)==float:
                       continue
                    clvalue="{}".format(clvalue2)


                    if clvalue in stopwordlist: # or clvalue is None or '1' in clvalue or '2' in clvalue or '3' in clvalue or '4' in clvalue or '5' in clvalue or '6' in clvalue or '7' in clvalue or '8' in clvalue or '9' in clvalue or '0' in clvalue:
                       continue
                    else:
                       clvalue3=removenew(clvalue)

                       for ant in clvalue3:
                           try:
                              int(ant)
                              eminem=True
                              continue
                           except:
                              eminem=False
                           if ant in xdct:
                               xdct[ant].append(dbse[entity])
                               continue
                           listn=[]
                           listn.append(dbse[entity])
                           xdct[ant]=listn
        else:
            continue
   thisfile2=open('C:\database1\DATABASEW.pkl','wb') #the stored file
   pickle.dump(xdct,thisfile2)
   thisfile2.close()
   thisfile1.close()

def mainu():
   for i in range(15):
      sleep(600)
      update_database()

threading.Thread(target=mainu).start()
